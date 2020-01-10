# python -m pip install python-pptx
import csv

from pptx import Presentation

# lines are interpreted thusly:
# layout, title, content/subtitle, notes
CONCEPT = "Concept"


def add_note_to_slide(slide, slide_notes):
    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = slide_notes


def parse_body_text(body_text):
    retval = body_text.replace("\\n", '\n')
    retval = retval.replace("\\r", '\r')
    retval = retval.replace("\\t", '\t')
    return retval


def add_background_to_slide(slide, slide_bkgnd):
    # retval = slide.shapes.add_picture(g, pic_left, pic_top, pic_width, pic_height)
    retval = slide.shapes.add_picture(slide_bkgnd, 0, 0)
    slide.shapes._spTree.remove(retval._element)
    slide.shapes._spTree.insert(2, retval._element)
    return retval


def get_placeholder_by_type(slide, ptype):
    retval = None
    for shape in slide.shapes:
        if shape.is_placeholder and shape.placeholder_format.type == ptype:
            retval = shape
            break
    return retval


def get_object_placeholder(slide):
    return get_placeholder_by_type(slide, 7)


def get_picture_placeholder(slide):
    return get_placeholder_by_type(slide, 18)


def add_image_to_slide(slide, slide_bkgnd):
    picture = None
    placeholder = get_picture_placeholder(slide)
    if (placeholder is not None):
        picture = placeholder.insert_picture(slide_bkgnd)
    else:
        picture = add_background_to_slide(slide,slide_bkgnd)

    return picture


def create_deck(titletxt, subtitletxt, name):
    prs = Presentation('template.pptx')
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = titletxt
    subtitle.text = subtitletxt

    with open(name + '.csv', 'rb') as f:
        reader = csv.reader(f)
        my_list = list(reader)

        for l in my_list:
            if len(l) == 0: continue

            slide_notes = None
            slide_bkgnd = None

            layout = l[0]
            if layout is None: layout = CONCEPT
            slide_title = l[1]
            slide_body = l[2]
            if len(l) > 3: slide_notes = l[3]
            if len(l) > 4: slide_bkgnd = l[4]

            blank_slide_layout = prs.slide_layouts.get_by_name(layout)
            slide = prs.slides.add_slide(blank_slide_layout)
            title = slide.shapes.title
            title.text = slide_title
            if slide_body is not None and slide_body > "":
                shapes = slide.shapes
                body_shape = shapes.placeholders[1]
                tf = body_shape.text_frame
                tf.text = parse_body_text(slide_body)

            #  if slide_bkgnd is not None and slide_bkgnd > "": add_background_to_slide(slide, slide_bkgnd)
            if slide_bkgnd is not None and slide_bkgnd > "": add_image_to_slide(slide, slide_bkgnd)
            if slide_notes is not None and slide_notes > "": add_note_to_slide(slide, slide_notes)

    prs.save(name + '.pptx')


if __name__ == "__main__":
    create_deck("Continuous Delivery Maturity Model", "Sustainable value delivery", "Overview")
    create_deck("Continuous Delivery Coaching", "Giving teams skills to deliver continually", "Coaching")
    create_deck("Culture & Organization", "The environment of continuous delivery", "Culture")
    create_deck("Design & Architecture", "Structuring your product for success", "Design")
    create_deck("Build & Deploy", "Building your product artifacts", "Build")
    create_deck("Testing & Verification", "Product artifact quality", "Testing")
    create_deck("Information & Reporting", "Measuring success and progress", "Information")
